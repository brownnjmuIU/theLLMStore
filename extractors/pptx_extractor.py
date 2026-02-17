from pptx import Presentation
import io


def extract_text_from_pptx(file_bytes: bytes) -> dict:
    """
    Extract text from PPTX bytes.
    Returns dict with text and slide_count.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    
    slides_text = []
    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_content.append(text)
        if slide_content:
            slides_text.append("\n".join(slide_content))
    
    full_text = "\n\n---\n\n".join(slides_text)  # separator between slides
    
    return {
        "text": full_text.strip(),
        "page_count": len(prs.slides)
    }
